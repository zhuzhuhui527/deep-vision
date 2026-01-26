#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = [
#     "openpyxl>=3.1.0",
#     "python-docx>=1.1.0",
#     "python-pptx>=0.6.23",
# ]
# ///
"""
Deep-Vision 文档格式转换工具

用途: 将 docx/xlsx/pptx 等格式转换为 Markdown 以便 AI 读取
使用方式: uvx scripts/convert_doc.py convert <文件路径> [输出目录]
"""

import argparse
import json
import logging
import os
import shutil
import sys
from datetime import datetime
from pathlib import Path
from typing import Optional

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s"
)
logger = logging.getLogger(__name__)

# 颜色代码（终端输出）
class Colors:
    RED = "\033[0;31m"
    GREEN = "\033[0;32m"
    YELLOW = "\033[1;33m"
    BLUE = "\033[0;34m"
    NC = "\033[0m"  # No Color


def log_info(message: str) -> None:
    """输出信息日志"""
    print(f"{Colors.GREEN}[INFO]{Colors.NC} {message}")


def log_warn(message: str) -> None:
    """输出警告日志"""
    print(f"{Colors.YELLOW}[WARN]{Colors.NC} {message}")


def log_error(message: str) -> None:
    """输出错误日志"""
    print(f"{Colors.RED}[ERROR]{Colors.NC} {message}")


def get_script_dir() -> Path:
    """获取脚本所在目录"""
    return Path(__file__).parent.resolve()


def setup_dirs(base_dir: Path) -> tuple[Path, Path]:
    """创建必要的目录"""
    temp_dir = base_dir / "data" / "temp"
    output_dir = base_dir / "data" / "converted"
    temp_dir.mkdir(parents=True, exist_ok=True)
    output_dir.mkdir(parents=True, exist_ok=True)
    return temp_dir, output_dir


def convert_docx_to_markdown(input_path: Path, output_path: Path) -> bool:
    """
    将 DOCX 文件转换为 Markdown

    Args:
        input_path: 输入 DOCX 文件路径
        output_path: 输出 Markdown 文件路径

    Returns:
        bool: 转换是否成功
    """
    try:
        from docx import Document

        log_info(f"转换 DOCX: {input_path}")
        doc = Document(str(input_path))

        markdown_lines = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                markdown_lines.append("")
                continue

            # 处理标题样式
            style_name = para.style.name.lower() if para.style else ""
            if "heading 1" in style_name:
                markdown_lines.append(f"# {text}")
            elif "heading 2" in style_name:
                markdown_lines.append(f"## {text}")
            elif "heading 3" in style_name:
                markdown_lines.append(f"### {text}")
            elif "heading 4" in style_name:
                markdown_lines.append(f"#### {text}")
            elif "list" in style_name:
                markdown_lines.append(f"- {text}")
            else:
                markdown_lines.append(text)

        # 处理表格
        for table in doc.tables:
            markdown_lines.append("")
            headers = []
            for i, row in enumerate(table.rows):
                cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
                row_text = "| " + " | ".join(cells) + " |"
                markdown_lines.append(row_text)
                if i == 0:
                    headers = cells
                    separator = "| " + " | ".join(["---"] * len(cells)) + " |"
                    markdown_lines.append(separator)
            markdown_lines.append("")

        output_path.write_text("\n".join(markdown_lines), encoding="utf-8")
        log_info(f"转换成功: {output_path}")
        return True

    except ImportError:
        log_error("缺少 python-docx 库，请使用 uvx 执行此脚本")
        return False
    except Exception as e:
        log_error(f"转换 DOCX 失败: {e}")
        return False


def convert_xlsx_to_markdown(input_path: Path, output_path: Path) -> bool:
    """
    将 XLSX 文件转换为 Markdown 表格

    Args:
        input_path: 输入 XLSX 文件路径
        output_path: 输出 Markdown 文件路径

    Returns:
        bool: 转换是否成功
    """
    try:
        from openpyxl import load_workbook

        log_info(f"转换 XLSX: {input_path}")
        wb = load_workbook(str(input_path), data_only=True)

        markdown_lines = [f"# {input_path.stem}", ""]

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            markdown_lines.append(f"## {sheet_name}")
            markdown_lines.append("")

            # 获取有数据的行和列
            rows = list(sheet.iter_rows(values_only=True))
            if not rows:
                markdown_lines.append("*空工作表*")
                continue

            # 过滤空行
            rows = [row for row in rows if any(cell is not None for cell in row)]
            if not rows:
                markdown_lines.append("*空工作表*")
                continue

            # 确定列数
            max_cols = max(len(row) for row in rows)

            for i, row in enumerate(rows):
                # 补齐列数
                cells = list(row) + [None] * (max_cols - len(row))
                cell_texts = [str(cell).replace("|", "\\|") if cell is not None else "" for cell in cells]
                row_text = "| " + " | ".join(cell_texts) + " |"
                markdown_lines.append(row_text)

                if i == 0:
                    separator = "| " + " | ".join(["---"] * max_cols) + " |"
                    markdown_lines.append(separator)

            markdown_lines.append("")

        output_path.write_text("\n".join(markdown_lines), encoding="utf-8")
        log_info(f"转换成功: {output_path}")
        return True

    except ImportError:
        log_error("缺少 openpyxl 库，请使用 uvx 执行此脚本")
        return False
    except Exception as e:
        log_error(f"转换 XLSX 失败: {e}")
        return False


def convert_pptx_to_markdown(input_path: Path, output_path: Path) -> bool:
    """
    将 PPTX 文件转换为 Markdown

    Args:
        input_path: 输入 PPTX 文件路径
        output_path: 输出 Markdown 文件路径

    Returns:
        bool: 转换是否成功
    """
    try:
        from pptx import Presentation

        log_info(f"转换 PPTX: {input_path}")
        prs = Presentation(str(input_path))

        markdown_lines = [f"# {input_path.stem}", ""]

        for slide_num, slide in enumerate(prs.slides, 1):
            markdown_lines.append(f"## 幻灯片 {slide_num}")
            markdown_lines.append("")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    # 判断是否为标题
                    if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                        if shape.placeholder_format.type == 1:  # 标题
                            markdown_lines.append(f"### {text}")
                        else:
                            markdown_lines.append(text)
                    else:
                        markdown_lines.append(text)
                    markdown_lines.append("")

                # 处理表格
                if shape.has_table:
                    table = shape.table
                    for i, row in enumerate(table.rows):
                        cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
                        row_text = "| " + " | ".join(cells) + " |"
                        markdown_lines.append(row_text)
                        if i == 0:
                            separator = "| " + " | ".join(["---"] * len(cells)) + " |"
                            markdown_lines.append(separator)
                    markdown_lines.append("")

            markdown_lines.append("---")
            markdown_lines.append("")

        output_path.write_text("\n".join(markdown_lines), encoding="utf-8")
        log_info(f"转换成功: {output_path}")
        return True

    except ImportError:
        log_error("缺少 python-pptx 库，请使用 uvx 执行此脚本")
        return False
    except Exception as e:
        log_error(f"转换 PPTX 失败: {e}")
        return False


def convert_document(input_path: str, output_dir: Optional[str] = None) -> Optional[str]:
    """
    转换单个文档

    Args:
        input_path: 输入文件路径
        output_dir: 输出目录（可选）

    Returns:
        Optional[str]: 转换后的文件路径，失败返回 None
    """
    input_file = Path(input_path).resolve()

    if not input_file.exists():
        log_error(f"文件不存在: {input_path}")
        return None

    script_dir = get_script_dir()
    base_dir = script_dir.parent
    _, default_output_dir = setup_dirs(base_dir)

    if output_dir:
        out_dir = Path(output_dir)
        out_dir.mkdir(parents=True, exist_ok=True)
    else:
        out_dir = default_output_dir

    ext = input_file.suffix.lower()
    output_path = out_dir / f"{input_file.stem}.md"

    converters = {
        ".docx": convert_docx_to_markdown,
        ".xlsx": convert_xlsx_to_markdown,
        ".xls": convert_xlsx_to_markdown,
        ".pptx": convert_pptx_to_markdown,
        ".ppt": convert_pptx_to_markdown,
    }

    if ext in [".md", ".txt"]:
        log_info("文件已是文本格式，直接复制")
        shutil.copy(input_file, output_path)
        return str(output_path)

    if ext == ".pdf":
        log_info("PDF 文件可直接由 Claude 读取，无需转换")
        return str(input_file)

    if ext not in converters:
        log_error(f"不支持的文件格式: {ext}")
        log_info("支持的格式: docx, xlsx, xls, pptx, ppt, md, txt, pdf")
        return None

    if converters[ext](input_file, output_path):
        return str(output_path)

    return None


def batch_convert(input_dir: str, output_dir: Optional[str] = None) -> dict:
    """
    批量转换目录中的文档

    Args:
        input_dir: 输入目录
        output_dir: 输出目录（可选）

    Returns:
        dict: 转换结果统计
    """
    input_path = Path(input_dir)

    if not input_path.is_dir():
        log_error(f"目录不存在: {input_dir}")
        return {"total": 0, "success": 0, "failed": 0}

    results = {"total": 0, "success": 0, "failed": 0, "files": []}

    for file_path in input_path.iterdir():
        if file_path.is_file():
            results["total"] += 1
            output = convert_document(str(file_path), output_dir)
            if output:
                results["success"] += 1
                results["files"].append({"input": str(file_path), "output": output})
            else:
                results["failed"] += 1

    log_info(f"转换完成: {results['success']}/{results['total']} 个文件成功")
    return results


def cleanup(base_dir: Optional[str] = None) -> None:
    """清理临时文件"""
    if base_dir:
        temp_dir = Path(base_dir) / "data" / "temp"
    else:
        temp_dir = get_script_dir().parent / "data" / "temp"

    if temp_dir.exists():
        shutil.rmtree(temp_dir)
        log_info("已清理临时文件")


def main():
    """主函数"""
    parser = argparse.ArgumentParser(
        description="Deep-Vision 文档格式转换工具",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  uvx scripts/convert_doc.py convert /path/to/document.docx
  uvx scripts/convert_doc.py batch /path/to/documents /path/to/output
  uvx scripts/convert_doc.py cleanup

支持的格式:
  输入: docx, xlsx, xls, pptx, ppt, md, txt, pdf
  输出: Markdown (.md)
        """
    )

    subparsers = parser.add_subparsers(dest="command", help="命令")

    # convert 命令
    convert_parser = subparsers.add_parser("convert", help="转换单个文件")
    convert_parser.add_argument("file", help="要转换的文件路径")
    convert_parser.add_argument("output", nargs="?", help="输出目录（可选）")

    # batch 命令
    batch_parser = subparsers.add_parser("batch", help="批量转换目录中的文件")
    batch_parser.add_argument("input_dir", help="输入目录")
    batch_parser.add_argument("output_dir", nargs="?", help="输出目录（可选）")

    # cleanup 命令
    subparsers.add_parser("cleanup", help="清理临时文件")

    # check 命令
    subparsers.add_parser("check", help="检查依赖是否可用")

    args = parser.parse_args()

    if args.command == "convert":
        result = convert_document(args.file, args.output)
        if result:
            print(result)
            sys.exit(0)
        else:
            sys.exit(1)

    elif args.command == "batch":
        results = batch_convert(args.input_dir, args.output_dir)
        print(json.dumps(results, ensure_ascii=False, indent=2))
        sys.exit(0 if results["failed"] == 0 else 1)

    elif args.command == "cleanup":
        cleanup()
        sys.exit(0)

    elif args.command == "check":
        try:
            from docx import Document
            from openpyxl import load_workbook
            from pptx import Presentation
            log_info("所有依赖已安装")
            sys.exit(0)
        except ImportError as e:
            log_error(f"缺少依赖: {e}")
            log_info("请使用 uvx 执行此脚本以自动安装依赖")
            sys.exit(1)
    else:
        parser.print_help()
        sys.exit(0)


if __name__ == "__main__":
    main()
